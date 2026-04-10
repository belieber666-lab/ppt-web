"""
PPTX 样式迁移 v3

修复：
- 模板可以没有 placeholder，从所有文本形状提取样式（按字号最大→标题，次大→正文）
- 支持 gradFill / solidFill / pattFill 颜色（之前只找 solidFill）
- 颜色搜索顺序：run.rPr → pPr.defRPr → lstStyle.lv1pPr.defRPr → lstStyle.defPPr.defRPr
- 如果任何层级都没有显式颜色，默认白色（防止深色背景上文字不可见）
"""

import copy
from typing import Optional
from lxml import etree

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE

# ── 命名空间 ──────────────────────────────────────────────────────────
A_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
IMG_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

FILL_TAGS = ("solidFill", "gradFill", "pattFill")

def _q(local: str) -> str:
    return f"{{{A_NS}}}{local}"

# ── 默认白色（兜底）────────────────────────────────────────────────────

def _white_fill_elem():
    """返回白色 solidFill XML 元素，作为颜色兜底。"""
    fill = etree.fromstring(
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="FFFFFF"/></a:solidFill>'
    )
    return fill


# ── 样式读取（从一个 rPr/defRPr 元素）─────────────────────────────────

def _read_elem(rPr, style: dict) -> None:
    """从 rPr / defRPr 元素读取尚未获取的字段写入 style。"""
    if rPr is None:
        return
    for attr in ("sz", "b", "i"):
        if style[attr] is None and rPr.get(attr) is not None:
            style[attr] = rPr.get(attr)
    if style["fill"] is None:
        for tag in FILL_TAGS:
            f = rPr.find(_q(tag))
            if f is not None:
                style["fill"] = copy.deepcopy(f)
                break
    if style["latin"] is None:
        lat = rPr.find(_q("latin"))
        if lat is not None:
            style["latin"] = copy.deepcopy(lat)


def _collect_style(text_frame) -> dict:
    """
    从文本框收集样式（XML 层面）。
    搜索顺序（由具体到抽象）：
      run.rPr → pPr.defRPr → lstStyle.lv1pPr.defRPr → lstStyle.defPPr.defRPr
    支持 solidFill / gradFill / pattFill。
    如果完全没有颜色，返回白色兜底。
    """
    style: dict = {
        "sz":    None,
        "b":     None,
        "i":     None,
        "fill":  None,   # <a:solidFill> 或 <a:gradFill> 副本
        "latin": None,   # <a:latin> 副本
    }

    # 1. 遍历所有段落
    for para in text_frame.paragraphs:
        # 段落级 defRPr（优先于 lstStyle，但低于 run）
        pPr = para._p.find(_q("pPr"))
        if pPr is not None:
            _read_elem(pPr.find(_q("defRPr")), style)
        # run 级（最高优先）
        for run in para.runs:
            _read_elem(run._r.find(_q("rPr")), style)

    # 2. 兜底：从 lstStyle 提取
    if style["fill"] is None or style["sz"] is None:
        txBody = text_frame._txBody
        lstStyle = txBody.find(_q("lstStyle"))
        if lstStyle is not None:
            # 检查 lvl1pPr … lvl9pPr（注意：是 lvl 不是 lv）
            for lvl in range(1, 10):
                lvpPr = lstStyle.find(_q(f"lvl{lvl}pPr"))
                if lvpPr is not None:
                    _read_elem(lvpPr.find(_q("defRPr")), style)
            # 检查 defPPr
            defPPr = lstStyle.find(_q("defPPr"))
            if defPPr is not None:
                _read_elem(defPPr.find(_q("defRPr")), style)

    # 3. 如果仍然没有颜色，使用白色（确保在深色背景上可见）
    if style["fill"] is None:
        style["fill"] = _white_fill_elem()

    return style


# ── 从模板提取所有文本形状样式，按字号分出"最大"和"次大"──────────────

def _get_shape_sz(text_frame) -> int:
    """获取文本框中最大的字号值（整数，0 表示未知）。"""
    max_sz = 0
    for para in text_frame.paragraphs:
        pPr = para._p.find(_q("pPr"))
        if pPr is not None:
            defRPr = pPr.find(_q("defRPr"))
            if defRPr is not None:
                sz = defRPr.get("sz")
                if sz:
                    max_sz = max(max_sz, int(sz))
        for run in para.runs:
            rPr = run._r.find(_q("rPr"))
            if rPr is not None:
                sz = rPr.get("sz")
                if sz:
                    max_sz = max(max_sz, int(sz))
        txBody = text_frame._txBody
        lstStyle = txBody.find(_q("lstStyle"))
        if lstStyle is not None:
            lv1 = lstStyle.find(_q("lvl1pPr"))
            if lv1 is not None:
                defRPr = lv1.find(_q("defRPr"))
                if defRPr is not None:
                    sz = defRPr.get("sz")
                    if sz:
                        max_sz = max(max_sz, int(sz))
    return max_sz


def extract_template_styles(template_slide) -> tuple:
    """
    从模板幻灯片提取标题和正文样式。
    策略：
      1. 先在 placeholder 中找（TITLE / BODY）
      2. 若 placeholder 不够，从所有文本形状按字号降序取 title / body
    返回 (title_style, body_style)
    """
    title_style: Optional[dict] = None
    body_style:  Optional[dict] = None

    # 优先从 placeholder 提取
    for shape in template_slide.shapes:
        if not shape.has_text_frame or not shape.is_placeholder:
            continue
        ph = shape.placeholder_format.type
        if ph in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            if title_style is None:
                title_style = _collect_style(shape.text_frame)
        else:
            if body_style is None:
                body_style = _collect_style(shape.text_frame)

    # 若 placeholder 未找到，从所有文本形状按字号排序
    if title_style is None or body_style is None:
        text_shapes = [
            (s, _get_shape_sz(s.text_frame))
            for s in template_slide.shapes
            if s.has_text_frame
        ]
        text_shapes.sort(key=lambda x: x[1], reverse=True)

        for shape, sz in text_shapes:
            style = _collect_style(shape.text_frame)
            if style["sz"] is None:
                style["sz"] = str(sz) if sz else None
            if title_style is None:
                title_style = style
                print(f"[pptx] 标题样式来自形状 '{shape.name}' sz={sz}")
            elif body_style is None and sz < (int(title_style.get("sz") or "9999")):
                body_style = style
                print(f"[pptx] 正文样式来自形状 '{shape.name}' sz={sz}")
            if title_style and body_style:
                break

    # 互为兜底
    if title_style is None:
        title_style = body_style or {"sz": None, "b": None, "i": None, "fill": _white_fill_elem(), "latin": None}
    if body_style is None:
        body_style = title_style

    return title_style, body_style


# ── 样式应用 ──────────────────────────────────────────────────────────

def _apply_style(text_frame, style: Optional[dict]) -> None:
    """将样式直接写入文本框每个 run 的 rPr（XML 层面）。"""
    if not style or all(v is None for v in style.values()):
        return

    for para in text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(_q("rPr"))
            if rPr is None:
                rPr = etree.Element(_q("rPr"))
                run._r.insert(0, rPr)

            # 字号
            if style["sz"] is not None:
                rPr.set("sz", style["sz"])
            # 粗体 / 斜体
            if style["b"] is not None:
                rPr.set("b", style["b"])
            if style["i"] is not None:
                rPr.set("i", style["i"])

            # 颜色（移除旧 fill，插入新 fill）
            if style["fill"] is not None:
                for tag in FILL_TAGS + ("noFill", "blipFill"):
                    old = rPr.find(_q(tag))
                    if old is not None:
                        rPr.remove(old)
                rPr.insert(0, copy.deepcopy(style["fill"]))

            # 字体
            if style["latin"] is not None:
                old_lat = rPr.find(_q("latin"))
                if old_lat is not None:
                    rPr.remove(old_lat)
                rPr.append(copy.deepcopy(style["latin"]))


# ── 背景复制 ──────────────────────────────────────────────────────────

def _copy_bg_fill(template_slide, target_slide) -> None:
    """复制 <p:bg> 元素（含图片关系重映射）到目标幻灯片。"""
    r_embed   = f"{{{R_NS}}}embed"
    tmpl_xml  = template_slide._element
    tgt_xml   = target_slide._element
    cSld_tmpl = tmpl_xml.find(f"{{{P_NS}}}cSld")
    cSld_tgt  = tgt_xml.find(f"{{{P_NS}}}cSld")
    if cSld_tmpl is None or cSld_tgt is None:
        return

    bg_tmpl = cSld_tmpl.find(f"{{{P_NS}}}bg")
    bg_tgt  = cSld_tgt.find(f"{{{P_NS}}}bg")
    if bg_tgt is not None:
        cSld_tgt.remove(bg_tgt)
    if bg_tmpl is None:
        return

    new_bg = copy.deepcopy(bg_tmpl)
    for elem in new_bg.iter():
        if r_embed in elem.attrib:
            old_rid = elem.attrib[r_embed]
            try:
                img_part = template_slide.part.related_part(old_rid)
                new_rid  = target_slide.part.relate_to(img_part, IMG_REL)
                elem.attrib[r_embed] = new_rid
            except Exception as exc:
                print(f"[warn] bg fill image: {exc}")
    cSld_tgt.insert(0, new_bg)


def _is_full_slide_pic(shape, w: int, h: int) -> bool:
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return False
        tol = int(w * 0.05)
        return (
            shape.left <= tol and shape.top <= tol
            and abs(shape.width - w) <= tol
            and abs(shape.height - h) <= tol
        )
    except Exception:
        return False


def _copy_bg_shapes(template_slide, target_slide, w: int, h: int) -> None:
    """复制模板中铺满全页的背景图片形状到目标幻灯片底层。"""
    r_embed = f"{{{R_NS}}}embed"
    bg_shapes = [s for s in template_slide.shapes if _is_full_slide_pic(s, w, h)]
    if not bg_shapes:
        return
    for s in list(target_slide.shapes):
        if _is_full_slide_pic(s, w, h):
            s._element.getparent().remove(s._element)
    spTree = target_slide.shapes._spTree
    for bg in reversed(bg_shapes):
        new_elem = copy.deepcopy(bg._element)
        for elem in new_elem.iter():
            if r_embed in elem.attrib:
                old_rid = elem.attrib[r_embed]
                try:
                    img_part = template_slide.part.related_part(old_rid)
                    new_rid  = target_slide.part.relate_to(img_part, IMG_REL)
                    elem.attrib[r_embed] = new_rid
                except Exception as exc:
                    print(f"[warn] bg shape image: {exc}")
        spTree.insert(2, new_elem)


# ── 主函数 ────────────────────────────────────────────────────────────

def process_pptx(template_path: str, content_path: str, output_path: str) -> None:
    """将模板第 1 页的视觉样式套用到内容 PPTX 每一页，输出到 output_path。"""
    tmpl_prs = Presentation(template_path)
    if not tmpl_prs.slides:
        raise ValueError("模板 PPT 没有幻灯片")

    tmpl_slide = tmpl_prs.slides[0]
    slide_w    = tmpl_prs.slide_width
    slide_h    = tmpl_prs.slide_height

    title_style, body_style = extract_template_styles(tmpl_slide)

    def _fmt(st):
        if st is None:
            return "None"
        fill_tag = st["fill"].tag.split("}")[-1] if st.get("fill") is not None else "None"
        return f"sz={st.get('sz')} fill={fill_tag} latin={st['latin'].get('typeface') if st.get('latin') is not None else None}"

    print(f"[pptx] 标题样式: {_fmt(title_style)}")
    print(f"[pptx] 正文样式: {_fmt(body_style)}")

    content_prs = Presentation(content_path)
    total = len(content_prs.slides)

    for idx, slide in enumerate(content_prs.slides):
        print(f"[pptx] 第 {idx + 1}/{total} 页...")
        _copy_bg_fill(tmpl_slide, slide)
        _copy_bg_shapes(tmpl_slide, slide, slide_w, slide_h)

        for shape in slide.shapes:
            if _is_full_slide_pic(shape, slide_w, slide_h):
                continue
            if shape.has_text_frame:
                is_title = (
                    shape.is_placeholder
                    and shape.placeholder_format.type in (
                        PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                    )
                )
                _apply_style(shape.text_frame, title_style if is_title else body_style)
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _apply_style(cell.text_frame, body_style)

    content_prs.save(output_path)
    print(f"[pptx] 完成 → {output_path}")
